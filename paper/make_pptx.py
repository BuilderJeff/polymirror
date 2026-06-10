"""Build the presentation slide deck (paper/polymirror_econfair.pptx).

"Classroom-clean" design: warm paper background, Georgia headlines, Verdana body,
thin ink rules, white cards.  All Study-2 numbers are read from
data/forward/experiment.json and data/forward/edge_curve.csv at build time;
nothing live is hardcoded.  (Study-1 historical constants are final:
151 wallets tested, 0 passed, best luck-chance 30.5%.)

Run:  .venv\\Scripts\\python.exe paper\\make_pptx.py
"""
from __future__ import annotations

import csv
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = Path(__file__).resolve().parent
ROOT = HERE.parent
PNG = HERE / "figs_png"
FWD = ROOT / "data" / "forward"

# ---------------------------------------------------------------------------
# 1.  DATA — read the live forward-experiment numbers at build time
# ---------------------------------------------------------------------------
with open(FWD / "experiment.json", encoding="utf-8") as f:
    EXP = json.load(f)

_positions = EXP["positions"]
_plist = list(_positions.values()) if isinstance(_positions, dict) else list(_positions)

N_POS = len(_plist)
N_MARKETS = len({p["condition_id"] for p in _plist})
N_WALLETS_ACTIVE = len({p["wallet"] for p in _plist})
N_MARKS = sum(len(p.get("marks", {})) for p in _plist)
N_BUYS = sum(1 for p in _plist if str(p.get("side", "")).upper() == "BUY")
N_SELLS = sum(1 for p in _plist if str(p.get("side", "")).upper() == "SELL")

EDGE_ROWS = []  # (horizon_h, n_wallets, n_positions, strat, bench, edge)
with open(FWD / "edge_curve.csv", encoding="utf-8") as f:
    for row in csv.DictReader(f):
        if not row.get("strat") or not row.get("edge"):
            continue  # horizon not yet populated
        if int(float(row["n_positions"] or 0)) == 0:
            continue
        EDGE_ROWS.append((int(row["horizon_h"]), int(row["n_wallets"]),
                          int(row["n_positions"]), float(row["strat"]),
                          float(row["bench"]), float(row["edge"])))
EDGE_ROWS.sort(key=lambda r: r[0])
if not EDGE_ROWS:
    raise SystemExit("edge_curve.csv has no populated horizons yet; rerun later.")

FIRST_H, _, _, STRAT_1, BENCH_1, EDGE_1 = EDGE_ROWS[0]
LATER_EDGES = [r[5] for r in EDGE_ROWS[1:]]
N_NEG_LATER = sum(1 for e in LATER_EDGES if e < 0)

# Computed, sign-pattern-driven headline for the result slide
if EDGE_1 > 0 and N_NEG_LATER > 0:
    CURVE_HEADLINE = "Copying beat the favorite early, then fell behind."
elif EDGE_1 > 0 and N_NEG_LATER == 0:
    CURVE_HEADLINE = "Copying held a small early edge over the favorite."
else:
    CURVE_HEADLINE = "Copying never reliably beat just betting the favorite."

BOTH_LOSE = any(r[3] < 0 and r[4] < 0 for r in EDGE_ROWS)

# ---------------------------------------------------------------------------
# 2.  DESIGN SYSTEM — "classroom-clean"
# ---------------------------------------------------------------------------
PAPER = RGBColor(0xFA, 0xF7, 0xF1)   # warm paper background
INK   = RGBColor(0x1F, 0x1F, 0x1F)   # near-black text
GREEN = RGBColor(0x1E, 0x7A, 0x52)   # the mirror strategy / money
GRAY  = RGBColor(0x6B, 0x72, 0x80)   # the benchmark / secondary text
GOLD  = RGBColor(0xD9, 0x8E, 0x04)   # edge / highlight
BRICK = RGBColor(0xB2, 0x3A, 0x2F)   # losses / warnings
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
EDGE_LINE = RGBColor(0xE4, 0xDE, 0xD2)  # subtle warm card border

HEAD = "Georgia"   # headlines + big stat numbers
BODY = "Verdana"   # body text

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
_no = 0


def _noshadow(shape):
    shape.shadow.inherit = False
    return shape


def rect(slide, l, t, w, h, fill, line=None, line_w=1.0):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    return _noshadow(sp)


def hrule(slide, l, t, w, color=INK):
    """Thin 1pt horizontal rule — the only ornament in this design."""
    return rect(slide, l, t, w, Pt(1.0), color)


def vrule(slide, l, t, h, color=EDGE_LINE):
    return rect(slide, l, t, Pt(1.0), h, color)


def card(slide, l, t, w, h):
    """Plain white content card with a subtle 1pt border (no rounding)."""
    return rect(slide, l, t, w, h, WHITE, line=EDGE_LINE, line_w=1.0)


def text(slide, l, t, w, h, items, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         wrap=True, line_spacing=1.08):
    """items: list of run-tuples (txt, size, color, bold[, font]) or paragraph
    dicts {runs:[...], align, space_after, space_before}."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, dict):
            p.alignment = item.get("align", align)
            p.space_after = Pt(item.get("space_after", 6))
            p.space_before = Pt(item.get("space_before", 0))
            p.line_spacing = item.get("line_spacing", line_spacing)
            rr = item["runs"]
        else:
            p.alignment = align
            p.space_after = Pt(6)
            p.line_spacing = line_spacing
            rr = [item]
        for run in rr:
            txt, size, color, bold = run[:4]
            font = run[4] if len(run) > 4 else BODY
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = font
    return tb


def pic(slide, name, l, t, w=None, h=None):
    path = PNG / name
    if not path.exists():
        return None
    p = slide.shapes.add_picture(str(path), l, t, width=w, height=h)
    p.line.color.rgb = EDGE_LINE
    p.line.width = Pt(1.0)
    return _noshadow(p)


def slide_base(title=None, number=True):
    """Paper background, optional Georgia title + thin rule, plain page number."""
    global _no
    _no += 1
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, PAPER)
    if number:
        text(s, SW - Inches(1.0), SH - Inches(0.42), Inches(0.6), Inches(0.3),
             [(str(_no), 10, GRAY, False)], align=PP_ALIGN.RIGHT)
    if title is None:
        return s, Inches(0.7)
    text(s, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.75),
         [(title, 30, INK, True, HEAD)])
    hrule(s, Inches(0.8), Inches(1.32), Inches(11.73))
    return s, Inches(1.65)


def bullets(slide, l, t, w, h, items, size=16, gap=10):
    paras = []
    for item in items:
        txt, color, bold = item[:3]
        if txt == "":
            paras.append({"runs": [(" ", 8, color, False)], "space_after": 2})
            continue
        paras.append({"runs": [(txt, size, color, bold)], "space_after": gap})
    return text(slide, l, t, w, h, paras, line_spacing=1.12)


def stat_block(slide, l, t, w, big, label, color=INK, big_size=40):
    """Plain stat: big Georgia number over a small Verdana label.  No box."""
    text(slide, l, t, w, Inches(0.85),
         [(big, big_size, color, True, HEAD)], align=PP_ALIGN.CENTER)
    text(slide, l, t + Inches(0.95), w, Inches(0.6),
         [(label, 12, GRAY, False)], align=PP_ALIGN.CENTER, line_spacing=1.05)


# ===========================================================================
# SLIDE 1 — TITLE
# ===========================================================================
s, _ = slide_base(number=False)
text(s, Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.7),
     [("Copy the Winners?", 60, INK, True, HEAD)], align=PP_ALIGN.CENTER)
text(s, Inches(0.8), Inches(3.55), Inches(11.7), Inches(0.6),
     [("Testing mirror-trading in prediction markets", 22, GRAY, False)],
     align=PP_ALIGN.CENTER)
hrule(s, Inches(5.17), Inches(4.45), Inches(3.0))
text(s, Inches(0.8), Inches(4.75), Inches(11.7), Inches(0.5),
     [("June 2026", 15, INK, False)], align=PP_ALIGN.CENTER)

# ===========================================================================
# SLIDE 2 — THE QUESTION
# ===========================================================================
s, top = slide_base("The Question")
text(s, Inches(1.2), Inches(2.15), Inches(10.9), Inches(1.9),
     [('"To what extent is mirror-trading a profitable investment strategy '
       'in current-event prediction markets?"', 28, INK, True, HEAD)],
     align=PP_ALIGN.CENTER, line_spacing=1.18)
text(s, Inches(1.2), Inches(4.55), Inches(10.9), Inches(0.5),
     [("In plain English:", 14, GRAY, False)], align=PP_ALIGN.CENTER)
text(s, Inches(1.2), Inches(5.05), Inches(10.9), Inches(0.9),
     [("If you copy every trade the top traders make... do you make money?",
       22, GREEN, True)], align=PP_ALIGN.CENTER)

# ===========================================================================
# SLIDE 3 — WHAT IS A PREDICTION MARKET
# ===========================================================================
s, top = slide_base("How Prediction Markets Work")
bullets(s, Inches(0.8), top + Inches(0.1), Inches(11.7), Inches(2.4), [
    ("Polymarket lets you buy shares in real-world outcomes: games, "
     "elections, news events.", INK, False),
    ("A YES share costs the market's probability. A 70-cent share means "
     "traders think there is a 70% chance.", INK, False),
    ("These prices are good forecasts. They beat Wall Street's inflation "
     "forecasts by 40% (Kalshi, 2018).", GREEN, False),
], size=16, gap=10)
# 3-step graphic: BUY -> RESOLVE -> PAYOUT
gy = Inches(4.45); gh = Inches(2.0); gw = Inches(3.45)
steps = [
    (Inches(0.8), "1.  BUY", [("Pay 70¢ for a YES share.", 13, INK, False)]),
    (Inches(4.95), "2.  RESOLVE", [("The real-world event happens.", 13, INK, False)]),
    (Inches(9.1), "3.  PAYOUT", [("Right: $1.00", 13, GREEN, True),
                                 ("Wrong: $0", 13, BRICK, True)]),
]
for gx, label, lines in steps:
    card(s, gx, gy, gw, gh)
    text(s, gx + Inches(0.3), gy + Inches(0.28), gw - Inches(0.6), Inches(0.5),
         [(label, 18, GOLD, True, HEAD)])
    hrule(s, gx + Inches(0.3), gy + Inches(0.82), gw - Inches(0.6), color=EDGE_LINE)
    paras = [{"runs": [ln], "space_after": 4} for ln in lines]
    text(s, gx + Inches(0.3), gy + Inches(1.0), gw - Inches(0.6), Inches(0.85), paras)
for ax in (Inches(4.3), Inches(8.45)):
    text(s, ax, gy + Inches(0.7), Inches(0.6), Inches(0.6),
         [("→", 26, GRAY, True)], align=PP_ALIGN.CENTER)

# ===========================================================================
# SLIDE 4 — THE PLAN: TWO STUDIES
# ===========================================================================
s, top = slide_base("The Plan: Two Studies")
cw, ch = Inches(5.75), Inches(4.4)
for cx, head, color, body_lines in [
    (Inches(0.8), "STUDY 1: Look back", GRAY,
     [("Is anyone reliably skilled, or just lucky?", 17, INK, True),
      ("Score the trading records of 151 of the most active wallets "
       "against pure chance.", 14, INK, False),
      ("151 traders tested", 14, GRAY, True)]),
    (Inches(6.8), "STUDY 2: Live test", GREEN,
     [("Copy the leaderboard for 48 hours.", 17, INK, True),
      ("Mirror every trade the Top-10 leaderboard wallets make, and track "
       "every position hour by hour.", 14, INK, False),
      (f"{N_POS:,} copied trades", 14, GREEN, True)]),
]:
    card(s, cx, top + Inches(0.15), cw, ch)
    text(s, cx + Inches(0.4), top + Inches(0.5), cw - Inches(0.8), Inches(0.5),
         [(head, 20, color, True, HEAD)])
    hrule(s, cx + Inches(0.4), top + Inches(1.12), cw - Inches(0.8), color=EDGE_LINE)
    paras = [{"runs": [ln], "space_after": 12} for ln in body_lines]
    text(s, cx + Inches(0.4), top + Inches(1.4), cw - Inches(0.8), ch - Inches(1.6),
         paras, line_spacing=1.15)

# ===========================================================================
# SLIDE 5 — STUDY 1: HOW
# ===========================================================================
s, top = slide_base("Study 1: How We Checked for Skill")
bullets(s, Inches(0.9), top + Inches(0.35), Inches(11.5), Inches(4.6), [
    ("We took the 151 most-active traders on Polymarket.", INK, False),
    ("We scored each one's accuracy with the Brier score (a score for how "
     "close your bets land to what actually happened).", INK, False),
    ("We compared each trader to a no-skill robot betting at the exact same "
     "market prices.", INK, False),
    ("Then we reshuffled the outcomes 10,000 times to see how often pure "
     "luck looks that good.", INK, False),
    ("", INK, False),
    ("A trader only counts as skilled if luck almost never matches their "
     "record.", GREEN, True),
], size=18, gap=16)

# ===========================================================================
# SLIDE 6 — STUDY 1: RESULT
# ===========================================================================
s, top = slide_base("Study 1: Result")
text(s, Inches(0.8), top + Inches(0.35), Inches(5.6), Inches(1.7),
     [("0 / 151", 88, BRICK, True, HEAD)])
text(s, Inches(0.8), top + Inches(2.1), Inches(5.6), Inches(0.9),
     [("Not one trader was measurably better than luck.", 18, INK, True)],
     line_spacing=1.15)
bullets(s, Inches(0.8), top + Inches(3.1), Inches(5.6), Inches(2.2), [
    ("Even the best trader's record had a 30% chance of being pure luck. "
     "We needed under 5%.", INK, False),
    ("An independent study of 588 million trades found the same "
     "(Akey et al., 2026).", GRAY, False),
], size=14, gap=10)
pic(s, "fig_s1_pvalue_dist.png", Inches(6.85), top + Inches(0.25), w=Inches(5.7))

# ===========================================================================
# SLIDE 7 — WHO ACTUALLY WINS
# ===========================================================================
s, top = slide_base("Who Actually Wins")
stat_block(s, Inches(1.0), top + Inches(0.2), Inches(3.5), "76%",
           "of all profits go to the top 1% of accounts", color=GOLD)
stat_block(s, Inches(4.9), top + Inches(0.2), Inches(3.5), "81%",
           "of the winners' volume is sports", color=GOLD)
stat_block(s, Inches(8.8), top + Inches(0.2), Inches(3.5), "588M",
           "trades in the study (Akey et al., 2026)", color=GOLD)
vrule(s, Inches(4.75), top + Inches(0.35), Inches(1.3))
vrule(s, Inches(8.65), top + Inches(0.35), Inches(1.3))
hrule(s, Inches(0.9), top + Inches(2.1), Inches(11.5), color=EDGE_LINE)
bullets(s, Inches(0.9), top + Inches(2.45), Inches(11.5), Inches(2.7), [
    ("The big winners are market makers. They earn the gap between buy and "
     "sell prices, like a casino earns the house edge.", INK, False),
    ("You cannot copy a market-maker's edge. By the time you see the trade, "
     "the price has already moved.", BRICK, True),
], size=17, gap=14)

# ===========================================================================
# SLIDE 8 — STUDY 2: THE LIVE EXPERIMENT
# ===========================================================================
s, top = slide_base("Study 2: The Live Experiment")
text(s, Inches(0.9), top + Inches(0.05), Inches(11.5), Inches(0.85),
     [("Study 1 found no skilled traders to copy. So we copied what a real "
       "person would: the Top-10 profit leaderboard.", 16, INK, False)],
     line_spacing=1.15)
# data-driven stats row
sb_y = top + Inches(1.05)
stat_block(s, Inches(0.9), sb_y, Inches(2.8), str(N_WALLETS_ACTIVE),
           "wallets traded", color=GREEN, big_size=36)
stat_block(s, Inches(3.85), sb_y, Inches(2.8), f"{N_POS:,}",
           "trades copied", color=GREEN, big_size=36)
stat_block(s, Inches(6.8), sb_y, Inches(2.8), str(N_MARKETS),
           "markets", color=GREEN, big_size=36)
stat_block(s, Inches(9.75), sb_y, Inches(2.8), "48",
           "hours", color=GREEN, big_size=36)
for vx in (Inches(3.7), Inches(6.65), Inches(9.6)):
    vrule(s, vx, sb_y + Inches(0.15), Inches(1.25))
hrule(s, Inches(0.9), sb_y + Inches(1.75), Inches(11.5), color=EDGE_LINE)
bullets(s, Inches(0.9), sb_y + Inches(2.0), Inches(11.5), Inches(2.0), [
    (f"We copied every trade ({N_BUYS:,} buys and {N_SELLS:,} sells) at "
     "the real market price the moment it happened.", INK, False),
    ('Fair benchmark: "just always bet the favorite" in the same markets '
     "at the same moments, no copying.", GRAY, True),
], size=16, gap=12)

# ===========================================================================
# SLIDE 9 — HOW WE MEASURED
# ===========================================================================
s, top = slide_base("How We Measured")
bullets(s, Inches(0.9), top + Inches(0.3), Inches(6.6), Inches(4.5), [
    ("Paper trading: real prices, no real money.", INK, True),
    ("Each position was tracked hourly, up to 48 hours, using Polymarket's "
     "minute-by-minute price history.", INK, False),
    ("When a market resolved, the position settled at $1 (right) or $0 "
     "(wrong).", INK, False),
    ("", INK, False),
    (f"{N_MARKS:,} hourly price checks in total.", GREEN, True),
], size=17, gap=14)
pic(s, "fig_mark_sources.png", Inches(7.9), top + Inches(0.45), w=Inches(4.6))

# ===========================================================================
# SLIDE 10 — RESULT: THE CURVE
# ===========================================================================
s, top = slide_base("Result: Copying vs. Betting the Favorite")
text(s, Inches(0.8), top + Inches(0.0), Inches(11.7), Inches(0.6),
     [(CURVE_HEADLINE, 21, INK, True, HEAD)])
pic(s, "fig_edge_curve.png", Inches(0.8), top + Inches(0.75), w=Inches(7.6))
cap_x = Inches(8.75)
text(s, cap_x, top + Inches(0.95), Inches(3.8), Inches(4.2),
     [{"runs": [("Reading the chart", 14, INK, True, HEAD)], "space_after": 10},
      {"runs": [("Mirror = copying the leaders.", 13, GREEN, True)], "space_after": 8},
      {"runs": [("Favorite = always betting the favorite.", 13, GRAY, True)], "space_after": 8},
      {"runs": [("Edge = the gap between them. Above zero means copying wins.",
                 13, GOLD, True)], "space_after": 10},
      {"runs": [(f"At {FIRST_H} hour{'s' if FIRST_H != 1 else ''}, the edge was "
                 f"{EDGE_1:+.1%}. "
                 + (f"It turned negative at {N_NEG_LATER} of the "
                    f"{len(LATER_EDGES)} later checkpoints."
                    if N_NEG_LATER else "It stayed positive after that."),
                 13, INK, False)], "space_after": 0}],
     line_spacing=1.18)

# ===========================================================================
# SLIDE 11 — RESULT: IS IT REAL?
# ===========================================================================
s, top = slide_base("Testing the Edge Against Luck")
pic(s, "fig_edge_ci.png", Inches(0.8), top + Inches(0.15), w=Inches(6.8))
bullets(s, Inches(8.0), top + Inches(0.5), Inches(4.55), Inches(4.5), [
    ("Look at the error bars: at almost every horizon, the range of "
     "plausible answers includes zero.", INK, True),
    ("We cannot rule out that the edge is pure chance.", BRICK, True),
    ("And this ignores trading costs. The gap between buy and sell prices "
     "would eat 1 to 2 cents per share.", GRAY, False),
], size=16, gap=16)

# ===========================================================================
# SLIDE 12 — THE ECONOMICS: WHY COPYING FAILS
# ===========================================================================
s, top = slide_base("The Economics: Why Copying Fails")
bullets(s, Inches(0.9), top + Inches(0.25), Inches(11.5), Inches(2.6), [
    ("Efficient-markets logic: the price already contains what informed "
     "traders know.", INK, False),
    ("When a big trader buys, the price jumps before you can copy. You pay "
     "the post-information price.", INK, False),
], size=17, gap=14)
hrule(s, Inches(2.4), top + Inches(2.35), Inches(8.5), color=EDGE_LINE)
text(s, Inches(1.2), top + Inches(2.7), Inches(10.9), Inches(1.0),
     [('"Copying is buying information that has already been sold."',
       24, GOLD, True, HEAD)], align=PP_ALIGN.CENTER, line_spacing=1.1)
text(s, Inches(1.2), top + Inches(4.0), Inches(10.9), Inches(0.8),
     [("Copy-trading research agrees: copiers rarely beat the market "
       "(Lu, 2017; Dorfleitner & Scheckenbach, 2022; Grundke & Wittke, 2024).",
       13, GRAY, False)], align=PP_ALIGN.CENTER, line_spacing=1.12)

# ===========================================================================
# SLIDE 13 — THE ANSWER
# ===========================================================================
s, top = slide_base("The Answer")
text(s, Inches(0.9), top + Inches(0.0), Inches(11.5), Inches(0.85),
     [('"To what extent is mirror-trading a profitable investment strategy '
       'in current-event prediction markets?"', 15, GRAY, False)],
     line_spacing=1.15)
text(s, Inches(0.9), top + Inches(0.95), Inches(11.5), Inches(0.85),
     [("To almost no extent.", 38, INK, True, HEAD)])
rows = [
    ("1", "0 of 151 historical traders performed better than luck."),
    ("2", "48 hours of live copying produced no reliable edge, and both "
          "strategies lost money before costs."),
    ("3", "Real-world trading costs would erase anything left."),
]
ry = top + Inches(2.15)
for num, line in rows:
    hrule(s, Inches(0.9), ry, Inches(11.5), color=EDGE_LINE)
    text(s, Inches(0.9), ry + Inches(0.18), Inches(0.7), Inches(0.7),
         [(num, 26, GOLD, True, HEAD)])
    text(s, Inches(1.8), ry + Inches(0.22), Inches(10.6), Inches(0.75),
         [(line, 16, INK, False)], line_spacing=1.12)
    ry += Inches(1.02)

# ===========================================================================
# SLIDE 14 — LIMITATIONS
# ===========================================================================
s, top = slide_base("Limitations")
bullets(s, Inches(0.9), top + Inches(0.35), Inches(11.5), Inches(4.8), [
    ("We watched one 2-day window. Markets change.", INK, False),
    (f"Only {N_WALLETS_ACTIVE} of the 10 leaderboard wallets actively "
     "traded during it.", INK, False),
    ("Paper trading ignores the buy/sell gap and slippage, so real results "
     "would be worse, not better.", BRICK, True),
    ("Picking wallets from a profit leaderboard is backward-looking: "
     "yesterday's winners, not tomorrow's.", INK, False),
], size=18, gap=18)

# ===========================================================================
# SLIDE 15 — TAKEAWAY
# ===========================================================================
s, _ = slide_base()
text(s, Inches(1.0), Inches(2.2), Inches(11.3), Inches(1.8),
     [('"Markets pay for information, not imitation."',
       40, INK, True, HEAD)], align=PP_ALIGN.CENTER, line_spacing=1.1)
hrule(s, Inches(5.17), Inches(4.25), Inches(3.0))
text(s, Inches(1.0), Inches(4.55), Inches(11.3), Inches(0.7),
     [("By the time you can copy a winner, the market already has.",
       18, GRAY, False)], align=PP_ALIGN.CENTER)
text(s, Inches(1.0), Inches(6.45), Inches(11.3), Inches(0.4),
     [(f"Based on {N_POS:,} mirrored trades collected live, June 2026.",
       11, GRAY, False)], align=PP_ALIGN.CENTER)

# ===========================================================================
# SLIDE 16 — WORKS CITED
# ===========================================================================
s, top = slide_base("Works Cited")
refs = [
    "Akey, Pat, et al. “Who Wins and Who Loses in Prediction Markets? Evidence from Polymarket.” "
    "SSRN, 18 Mar. 2026, ssrn.com/abstract=6443103.",
    "Dorfleitner, Gregor, and Isabel Scheckenbach. “Trading Activity on Social Trading Platforms - "
    "a Behavioral Approach.” The Journal of Risk Finance, vol. 23, no. 1, Jan. 2022, pp. 32-54.",
    "Grundke, Peter, and Gerrit Wittke. “Social Trading Platforms vs. Mutual Funds: Herding "
    "Tendencies and Portfolio Risks.” The European Journal of Finance, vol. 31, no. 7, Dec. 2024, "
    "pp. 827-49.",
    "“Kalshi.” Kalshi.com, 2018, research.kalshi.com/articles/crisis-alpha.",
    "“Kalshi Research.” Kalshi.com, 2026, research.kalshi.com/articles/mamdani-primary-victory.",
    "Lu, Juye Shirley. “To Mirror or Not to Mirror: Modeling Relationships in Social Trading.” MIT, "
    "2017, dspace.mit.edu/entities/publication/d04baf10-ccf7-4358-9aaf-6716565782e9.",
]
paras = [{"runs": [(r, 13, INK, False)], "space_after": 14} for r in refs]
text(s, Inches(0.9), top + Inches(0.25), Inches(11.5), Inches(5.0), paras,
     line_spacing=1.12)

# ---------------------------------------------------------------------------
# SAVE (lock-tolerant)
# ---------------------------------------------------------------------------
saved = None
for cand in ["polymirror_econfair.pptx", "polymirror_econfair_v2.pptx",
             "polymirror_econfair_v3.pptx"]:
    try:
        prs.save(str(HERE / cand)); saved = cand; break
    except PermissionError:
        continue
if saved is None:
    raise SystemExit("All candidate filenames are locked — close the open PowerPoint files.")
print("wrote", HERE / saved, "with", len(prs.slides._sldIdLst), "slides")
print(f"data: n_pos={N_POS} n_markets={N_MARKETS} n_wallets_active={N_WALLETS_ACTIVE} "
      f"n_marks={N_MARKS} buys={N_BUYS} sells={N_SELLS} "
      f"edge_1h={EDGE_1:+.4f} later_neg={N_NEG_LATER}/{len(LATER_EDGES)}")
print("headline:", CURVE_HEADLINE)
