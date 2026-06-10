"""Build the symposium slide deck (polymirror_symposium.pptx) — designed edition.
Run after the figure generators:  .venv\\Scripts\\python.exe paper\\make_pptx.py
Pure-local; depends only on python-pptx + the PNGs in paper/figs_png/.
"""
from __future__ import annotations
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = Path(__file__).resolve().parent
PNG = HERE / "figs_png"

# ---- palette (shared with the figures) -------------------------------------
INK    = RGBColor(0x10, 0x2A, 0x43)   # near-black navy — backgrounds, titles
INK2   = RGBColor(0x16, 0x35, 0x52)   # slightly lighter navy band
BLUE   = RGBColor(0x2C, 0x6F, 0xBB)   # primary
TEAL   = RGBColor(0x1A, 0xA3, 0x9A)   # accent
AMBER  = RGBColor(0xE8, 0xA3, 0x3D)   # highlight
RED    = RGBColor(0xD6, 0x45, 0x50)   # caution
SLATE  = RGBColor(0x53, 0x60, 0x70)   # secondary text
MUTE   = RGBColor(0x8A, 0x97, 0xA6)   # tertiary
CARD   = RGBColor(0xEE, 0xF3, 0xF8)   # light card fill
CARD2  = RGBColor(0xF6, 0xF9, 0xFC)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
CLOUD  = RGBColor(0xC9, 0xD6, 0xE5)   # light text on navy
HEAD   = "Calibri"
BODY   = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
_no = 0


# ---- low-level helpers -----------------------------------------------------
def _noshadow(shape):
    shape.shadow.inherit = False
    return shape


def rect(slide, l, t, w, h, color, line=None, line_w=1.0):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    return _noshadow(sp)


def rrect(slide, l, t, w, h, color, radius=0.08, line=None, line_w=1.25):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    return _noshadow(sp)


def tri(slide, l, t, w, h, color, flip=False):
    sp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, l, t, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    if flip:
        sp.rotation = 180
    return _noshadow(sp)


def text(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         wrap=True, line_spacing=1.0):
    """runs: list of (text, size, color, bold) OR list of paragraphs where each paragraph
    is a dict {runs:[...], align, space_after, level, bullet}."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, item in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, dict):
            p.alignment = item.get("align", align)
            p.space_after = Pt(item.get("space_after", 6))
            p.space_before = Pt(item.get("space_before", 0))
            p.level = item.get("level", 0)
            if line_spacing:
                p.line_spacing = line_spacing
            pre = item.get("bullet", "")
            rr = item["runs"]
            for j, (txt, size, color, bold) in enumerate(rr):
                r = p.add_run(); r.text = (pre if j == 0 else "") + txt
                r.font.size = Pt(size); r.font.bold = bold
                r.font.color.rgb = color; r.font.name = BODY
        else:
            txt, size, color, bold = item
            p.alignment = align
            p.space_after = Pt(6)
            if line_spacing:
                p.line_spacing = line_spacing
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = BODY
    return tb


def pic(slide, name, l, t, w=None, h=None, frame=True):
    path = PNG / name
    if not path.exists():
        return None
    p = slide.shapes.add_picture(str(path), l, t, width=w, height=h)
    if frame:
        p.line.color.rgb = RGBColor(0xDD, 0xE4, 0xEC)
        p.line.width = Pt(1.0)
    _noshadow(p)
    return p


# ---- slide chrome ----------------------------------------------------------
def content(title, kicker):
    global _no
    _no += 1
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, WHITE)
    rect(s, 0, 0, Inches(0.16), SH, TEAL)            # left accent strip
    text(s, Inches(0.55), Inches(0.42), Inches(11.5), Inches(0.3),
         [(kicker.upper(), 12.5, TEAL, True)])
    text(s, Inches(0.55), Inches(0.70), Inches(12.2), Inches(0.8),
         [(title, 27, INK, True)], anchor=MSO_ANCHOR.TOP)
    rect(s, Inches(0.55), Inches(1.52), Inches(12.2), Pt(1.6), RGBColor(0xE2, 0xE8, 0xEF))
    # footer
    rect(s, Inches(0.55), SH - Inches(0.52), Inches(12.2), Pt(1.2), RGBColor(0xE2, 0xE8, 0xEF))
    text(s, Inches(0.55), SH - Inches(0.45), Inches(10.5), Inches(0.3),
         [("Mirror-trading on Polymarket  ·  leakage-controlled backtest + live forward test",
           9.5, MUTE, False)])
    text(s, SW - Inches(1.1), SH - Inches(0.47), Inches(0.55), Inches(0.32),
         [(f"{_no:02d}", 11, TEAL, True)], align=PP_ALIGN.RIGHT)
    return s, Inches(1.78)


def divider(num, title, sub, accent):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, INK)
    rect(s, 0, 0, SW, SH, INK)
    # accent geometry
    rect(s, 0, Inches(2.55), Inches(2.2), Pt(4), accent)
    tri(s, SW - Inches(2.6), 0, Inches(2.6), Inches(2.6), INK2)
    rrect(s, SW - Inches(1.9), Inches(0.55), Inches(1.25), Inches(1.25), accent, radius=0.22)
    text(s, SW - Inches(1.9), Inches(0.55), Inches(1.25), Inches(1.25),
         [(f"{num:02d}", 40, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(0.9), Inches(2.75), Inches(11.0), Inches(1.3),
         [(title, 40, WHITE, True)])
    text(s, Inches(0.92), Inches(4.05), Inches(10.5), Inches(1.2),
         [(sub, 18, CLOUD, False)], line_spacing=1.15)
    return s


def stat_card(slide, l, t, w, h, big, label, accent=TEAL, big_size=30, fill=CARD):
    rrect(slide, l, t, w, h, fill, radius=0.10)
    rect(slide, l, t, Inches(0.10), h, accent)       # accent edge
    text(slide, l + Inches(0.28), t + Inches(0.16), w - Inches(0.4), h - Inches(0.3),
         [{"runs": [(big, big_size, accent, True)], "space_after": 2},
          {"runs": [(label, 12.5, SLATE, False)], "space_after": 0}],
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)


def bullets(slide, l, t, w, h, items, size=16, gap=8):
    paras = []
    for (txt, lvl, color, bold) in items:
        if txt == "":
            paras.append({"runs": [(" ", 6, color, False)], "space_after": 2})
            continue
        bullet = "▪  " if lvl == 0 else "–  "
        paras.append({"runs": [(bullet + txt, size - lvl, color, bold)],
                      "level": lvl, "space_after": gap})
    return text(slide, l, t, w, h, paras, line_spacing=1.06)


def chip(slide, l, t, w, h, code, desc):
    rrect(slide, l, t, w, h, CARD, radius=0.18, line=RGBColor(0xD7, 0xE0, 0xEA), line_w=1.0)
    rrect(slide, l + Inches(0.10), t + Inches(0.12), Inches(0.62), h - Inches(0.24), INK, radius=0.30)
    text(slide, l + Inches(0.10), t + Inches(0.12), Inches(0.62), h - Inches(0.24),
         [(code, 12.5, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(slide, l + Inches(0.82), t, w - Inches(0.92), h,
         [(desc, 11.5, SLATE, False)], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)


# =====================================================================
# 1 — TITLE
# =====================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, INK)
rect(s, 0, 0, Inches(4.6), SH, INK2)                       # left band
rect(s, Inches(4.6), 0, Pt(4), SH, TEAL)                   # divider rule
# corner accents
rrect(s, Inches(0.7), Inches(0.7), Inches(0.95), Inches(0.95), TEAL, radius=0.24)
text(s, Inches(0.7), Inches(0.7), Inches(0.95), Inches(0.95),
     [("PM", 26, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, Inches(0.7), Inches(1.9), Inches(3.4), Inches(0.4),
     [("POLYMIRROR", 15, TEAL, True)])
text(s, Inches(0.7), Inches(5.7), Inches(3.5), Inches(1.4),
     [{"runs": [("Read-only study of", 12.5, CLOUD, False)], "space_after": 2},
      {"runs": [("prediction-market", 12.5, CLOUD, False)], "space_after": 2},
      {"runs": [("mirror-trading", 12.5, CLOUD, False)], "space_after": 0}], line_spacing=1.1)
# right: title
text(s, Inches(5.1), Inches(1.55), Inches(7.6), Inches(2.6),
     [("Does Copying Skilled Wallets Beat Buying the Favorite?", 38, WHITE, True)],
     line_spacing=1.02)
text(s, Inches(5.1), Inches(4.35), Inches(7.5), Inches(1.2),
     [("Evidence from a leakage-controlled backtest and a live forward test on Polymarket",
       19, CLOUD, False)], line_spacing=1.12)
rect(s, Inches(5.12), Inches(5.75), Inches(2.0), Pt(3), TEAL)
text(s, Inches(5.1), Inches(6.05), Inches(7.5), Inches(0.9),
     [{"runs": [("polymirror project", 13, WHITE, True)], "space_after": 1},
      {"runs": [("Run dated 2026-06-07/08  ·  all code, parameters, and seeds in the repository",
                 11, MUTE, False)], "space_after": 0}])

# =====================================================================
# 2 — DIVIDER 01
# =====================================================================
divider(1, "Setup", "The question, the rules of evidence, and the data.", TEAL)

# =====================================================================
# 3 — THE QUESTION + R-RULES
# =====================================================================
s, top = content("The question — and the standard of proof", "Motivation")
bullets(s, Inches(0.55), top, Inches(7.0), Inches(4.9), [
    ("Can mirroring “good” Polymarket wallets beat buying the favorite, after costs, at short "
     "horizons?", 0, INK, True),
    ("Why care? Prediction-market prices beat Wall St. consensus on CPI by 40.1% MAE "
     "(Kalshi, 2018).", 0, BLUE, False),
    ("Unit of analysis = the on-chain wallet. A wallet is not a person.", 0, SLATE, False),
    ("Headline = an edge-decay curve (mirror − favorite) vs. horizon — not one number.", 0, SLATE, False),
    ("Standard: empirical honesty, not profit. A defensible null is a success.", 0, SLATE, False),
    ("Two independent attacks:", 0, INK, True),
    ("Study 1 — strict historical skill-vs-luck filter (large sample).", 1, BLUE, True),
    ("Study 2 — live forward mirror experiment (real time).", 1, TEAL, True),
], size=16)
# right column: R-rule chips
text(s, Inches(8.0), top - Inches(0.05), Inches(4.8), Inches(0.35),
     [("DESIGN RULES — VIOLATE ONE, INVALIDATE THE RESULT", 10.5, MUTE, True)])
chips = [("R0", "Read-only; no live trade or key"),
         ("R1", "Strict temporal split — no leakage"),
         ("R2", "Select by proper score, never PnL"),
         ("R3", "Must beat a luck null, not just rank"),
         ("R4", "Favorite fixed at entry, no hindsight"),
         ("R5", "Strategy = benchmark except the side"),
         ("R6", "Spread swept over three presets"),
         ("R8", "One config; every step seeded")]
cy = top + Inches(0.32)
for code, desc in chips:
    chip(s, Inches(8.0), cy, Inches(4.8), Inches(0.52), code, desc)
    cy += Inches(0.595)

# =====================================================================
# 4 — DATA & APIS
# =====================================================================
s, top = content("Data: three public, read-only APIs", "Sources")
bullets(s, Inches(0.55), top, Inches(7.1), Inches(4.8), [
    ("Gamma — markets & ground-truth resolution.", 0, INK, True),
    ("Data API — per-wallet / per-market trades.", 0, INK, True),
    ("/activity honors start/end (forward capture).", 1, SLATE, False),
    ("/trades capped at ~4,000 recent records.", 1, SLATE, False),
    ("CLOB — /midpoint, /spread, /prices-history.", 0, INK, True),
], size=16)
# correction callout card
rrect(s, Inches(8.0), top, Inches(4.8), Inches(4.3), CARD, radius=0.08)
rect(s, Inches(8.0), top, Inches(0.12), Inches(4.3), AMBER)
text(s, Inches(8.28), top + Inches(0.25), Inches(4.35), Inches(3.9),
     [{"runs": [("METHODOLOGICAL CORRECTION", 11.5, AMBER, True)], "space_after": 8},
      {"runs": [("/prices-history returns a 1-minute mid series that ", 13.5, INK, False),
                ("persists for days after resolution", 13.5, INK, True),
                (".", 13.5, INK, False)], "space_after": 8},
      {"runs": [("Verified: a market resolved ~6 days earlier still returned its full history.",
                 12, SLATE, False)], "space_after": 10},
      {"runs": [("→ exit marks reconstructed retrospectively — the collector needs the network "
                 "only at capture/analysis, not at every horizon.", 13, TEAL, True)], "space_after": 0}],
     line_spacing=1.08)

# =====================================================================
# 5 — DIVIDER 02
# =====================================================================
divider(2, "Study 1 — Historical", "Does any wallet show calibration skill beyond the market "
        "price? A leakage-controlled luck filter over 151 wallets.", BLUE)

# =====================================================================
# 6 — STUDY 1 METHOD
# =====================================================================
s, top = content("Skill-vs-luck filter", "Study 1 · Method")
bullets(s, Inches(0.55), top, Inches(12.2), Inches(5.0), [
    ("Collapse all fills to ONE position per (wallet, market, outcome) — a 71-fill market is one "
     "correlated bet, not 71. This kills inflated naive t-stats.", 0, INK, False),
    ("Score each wallet by the Brier proper score on entry price vs. realized outcome.", 0, INK, False),
    ("Luck null: hold prices pᵢ fixed, redraw outcomes yᵢ* ~ Bernoulli(pᵢ) 10,000× — the "
     "distribution of NO skill beyond the market price.", 0, INK, False),
    ("Laplace-smoothed empirical p-value · Benjamini–Hochberg FDR across wallets · per-wallet "
     "seeded RNG.", 0, INK, False),
    ("", 0, SLATE, False),
    ("The null is deliberately strict — buying at 0.84 and winning 84% of the time is exactly what "
     "it does, and earns no significance (it subsumes the favorite-buyer fallacy).", 0, TEAL, True),
], size=16, gap=11)

# =====================================================================
# 7 — STUDY 1 RESULT
# =====================================================================
s, top = content("Result: zero wallets survive", "Study 1 · Result")
stat_card(s, Inches(0.55), top, Inches(2.95), Inches(1.35), "0 / 151", "wallets survive (raw & FDR)",
          accent=RED, big_size=28)
stat_card(s, Inches(0.55), top + Inches(1.55), Inches(2.95), Inches(1.35), "0.305",
          "minimum p-value (none < 0.30)", accent=BLUE, big_size=30)
stat_card(s, Inches(0.55), top + Inches(3.10), Inches(2.95), Inches(1.35), "~7–8",
          "raw winners expected by chance", accent=AMBER, big_size=30)
pic(s, "fig_s1_pvalue_dist.png", Inches(3.95), top + Inches(0.1), h=Inches(4.15))
text(s, Inches(3.85), top + Inches(4.45), Inches(8.9), Inches(0.4),
     [("Under R3, the data admit no wallet worth mirroring — echoing Akey et al. (2026).",
       12.5, TEAL, True)], align=PP_ALIGN.CENTER)

# =====================================================================
# 8 — STUDY 1 EVIDENCE
# =====================================================================
s, top = content("The population shows no skill", "Study 1 · Evidence")
pic(s, "fig_s1_brier_scatter.png", Inches(0.7), top + Inches(0.05), h=Inches(4.55))
pic(s, "fig_s1_null_margin.png", Inches(6.35), top + Inches(0.6), w=Inches(6.4))
text(s, Inches(0.7), top + Inches(4.75), Inches(12.0), Inches(0.4),
     [("Observed Brier ≥ price-null Brier for most wallets (left); null margin centered near zero, "
       "mean +0.0024 (right).", 11.5, SLATE, False)], align=PP_ALIGN.CENTER)

# =====================================================================
# 9 — DIVIDER 03
# =====================================================================
divider(3, "Study 2 — Forward", "An empty watchlist forces an exploratory pivot: mirror the live "
        "leaderboard and mark every trade in real time.", TEAL)

# =====================================================================
# 10 — STUDY 2 SELECTION
# =====================================================================
s, top = content("A live forward test (explicit R2 deviation)", "Study 2 · Design")
bullets(s, Inches(0.55), top, Inches(7.1), Inches(4.8), [
    ("Study 1 left an empty watchlist — a confirmatory “skilled-wallet” test was impossible.",
     0, INK, False),
    ("So we ran an exploratory test with a different, openly-labeled rule:", 0, INK, False),
    ("mirror the top-10 wallets by 30-day leaderboard profit, each verified active.", 1, TEAL, True),
    ("This violates R2 (selection by PnL) — hypothesis-generating, not confirmatory.", 1, RED, True),
], size=16, gap=10)
rrect(s, Inches(8.0), top, Inches(4.8), Inches(3.0), INK, radius=0.08)
text(s, Inches(8.3), top + Inches(0.3), Inches(4.3), Inches(2.5),
     [{"runs": [("THE PRACTICAL QUESTION", 11.5, TEAL, True)], "space_after": 10},
      {"runs": [("If a retail user naively copies whoever is winning right now, do they beat "
                 "buying the favorite?", 18, WHITE, True)], "space_after": 0}], line_spacing=1.12)

# =====================================================================
# 11 — STUDY 2 METHOD
# =====================================================================
s, top = content("Capture & retrospective marking", "Study 2 · Method")
bullets(s, Inches(0.55), top, Inches(12.2), Inches(5.0), [
    ("Window 20:04 → 08:04 UTC. Capture EVERY trade (buy & sell) in live binary markets.", 0, INK, False),
    ("Normalize each to an opening long (sell of k @ p ≡ long of complement @ 1−p): copy the action.",
     0, INK, False),
    ("Reconstruct both token mids at entry and each elapsed hour N from /prices-history; horizons "
     "past resolution settle to 0/1.", 0, INK, False),
    ("Strategy leg = copied long; benchmark = favorite-at-entry. Identical except the side (R5).",
     0, INK, False),
    ("Collapse to one position per (wallet, market, side), then cluster by wallet.", 0, INK, False),
    ("", 0, SLATE, False),
    ("All 23,584 marks reconstructed — zero missing. Laptop offline between check-ins loses nothing.",
     0, TEAL, True),
], size=16, gap=9)

# =====================================================================
# 12 — STUDY 2 SAMPLE
# =====================================================================
s, top = content("The sample — 5,854 trades, 202 markets", "Study 2 · Sample")
pic(s, "fig_wallets.png", Inches(0.6), top + Inches(0.15), w=Inches(4.0))
pic(s, "fig_entry_price.png", Inches(4.75), top + Inches(0.15), w=Inches(4.0))
pic(s, "fig_fills_per_market.png", Inches(8.9), top + Inches(0.15), w=Inches(4.0))
text(s, Inches(0.6), top + Inches(3.5), Inches(12.3), Inches(0.7),
     [("4 of 10 wallets traded (inference rests on 3)  ·  entry prices span the unit interval "
       "(median 0.48)  ·  fills per market heavy-tailed (max 1,167 → dedup).", 12, SLATE, False)],
     align=PP_ALIGN.CENTER, line_spacing=1.1)

# =====================================================================
# 13 — STUDY 2 RESULT
# =====================================================================
s, top = content("A positive edge — but fragile", "Study 2 · Result")
pic(s, "fig_edge_curve.png", Inches(0.55), top + Inches(0.1), w=Inches(6.05))
pic(s, "fig_edge_ci.png", Inches(6.8), top + Inches(0.1), w=Inches(6.05))
text(s, Inches(0.55), top + Inches(3.7), Inches(12.3), Inches(0.7),
     [{"runs": [("Edge rises with horizon (+2.0% @1h → +12.6% @7h) — but the by-market 95% CI "
                 "includes zero at ", 12, SLATE, False),
                ("every", 12, RED, True),
                (" horizon, and both legs lose money gross at ≥4h.", 12, SLATE, False)]}],
     align=PP_ALIGN.CENTER, line_spacing=1.1)

# =====================================================================
# 14 — WHY NOT REAL
# =====================================================================
s, top = content("Why the edge is not a finding", "Study 2 · Interpretation")
cards = [
    ("1", "Clustering changes everything",
     "By wallet (n=4) it looks significant (p=0.006) — but that's just four wallets sharing a sign. "
     "By market (n≈200–340) every 95% CI includes zero (min p=0.061)."),
    ("2", "Cohort-dependent & unstable",
     "The long-horizon edge exists only for the afternoon cohort. When fresh trades entered the "
     "short horizons, 1–3h edge collapsed and 3h flipped NEGATIVE."),
    ("3", "Both legs lose money",
     "At ≥4h both mirror and favorite return −15% to −28% gross. “Mirror loses less,” not "
     "“mirror profits.” And this is BEFORE spread costs."),
]
cw = Inches(3.95); cx = Inches(0.55); gap = Inches(0.27)
for nbox, title, body in cards:
    rrect(s, cx, top + Inches(0.1), cw, Inches(4.3), CARD, radius=0.07)
    rect(s, cx, top + Inches(0.1), cw, Inches(0.12), RED)
    rrect(s, cx + Inches(0.28), top + Inches(0.42), Inches(0.7), Inches(0.7), RED, radius=0.30)
    text(s, cx + Inches(0.28), top + Inches(0.42), Inches(0.7), Inches(0.7),
         [(nbox, 24, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, cx + Inches(0.3), top + Inches(1.4), cw - Inches(0.6), Inches(0.9),
         [(title, 16.5, INK, True)], line_spacing=1.0)
    text(s, cx + Inches(0.3), top + Inches(2.3), cw - Inches(0.6), Inches(1.9),
         [(body, 12.5, SLATE, False)], line_spacing=1.12)
    cx += cw + gap

# =====================================================================
# 15 — DIVIDER 04
# =====================================================================
divider(4, "Synthesis", "Two independent methods converge on one answer.", AMBER)

# =====================================================================
# 16 — DISCUSSION
# =====================================================================
s, top = content("Two methods, one answer", "Discussion")
# two side cards
rrect(s, Inches(0.55), top, Inches(6.0), Inches(2.0), CARD, radius=0.08)
rect(s, Inches(0.55), top, Inches(0.12), Inches(2.0), BLUE)
text(s, Inches(0.85), top + Inches(0.22), Inches(5.5), Inches(1.7),
     [{"runs": [("STUDY 1 — rigorous, large", 11.5, BLUE, True)], "space_after": 6},
      {"runs": [("No wallet shows calibration skill beyond the market price.", 15, INK, False)],
       "space_after": 0}], line_spacing=1.1)
rrect(s, Inches(6.75), top, Inches(6.0), Inches(2.0), CARD, radius=0.08)
rect(s, Inches(6.75), top, Inches(0.12), Inches(2.0), TEAL)
text(s, Inches(7.05), top + Inches(0.22), Inches(5.5), Inches(1.7),
     [{"runs": [("STUDY 2 — relaxed, practical", 11.5, TEAL, True)], "space_after": 6},
      {"runs": [("Copying current winners gives a fragile, insignificant, gross-of-cost edge over "
                 "an also-losing benchmark.", 15, INK, False)], "space_after": 0}], line_spacing=1.1)
rrect(s, Inches(0.55), top + Inches(2.25), Inches(12.2), Inches(2.5), INK, radius=0.06)
text(s, Inches(0.95), top + Inches(2.55), Inches(11.4), Inches(2.0),
     [{"runs": [("Convergence = evidence consistent with short-horizon efficiency", 19, WHITE, True)],
       "space_after": 8},
      {"runs": [("By the time a trade is observable, its information is already in the price.",
                 15, CLOUD, False)], "space_after": 8},
      {"runs": [("Inference is one-directional: we failed to reject efficiency under two reasonable "
                 "strategies — we do not claim to have proven it.", 14, TEAL, True)], "space_after": 0}],
     line_spacing=1.1)

# =====================================================================
# 16b — RELATED WORK / OUR RESULT IN CONTEXT
# =====================================================================
s, top = content("Our result in the literature", "Related work")
lit = [
    ("Akey et al. (2026)\nSSRN · Polymarket",
     "588M trades, $67B: the top 1% of traders earn 76% of profits, and the winners are "
     "market makers in narrow markets (81% sports) — not copyable speculators.",
     "Corroborates our zero skilled-speculator survivors & sports-heavy sample.", BLUE),
    ("Lu (2017)\nMIT · eToro, 87.5M trades",
     "Copying expert traders reduces risk but rarely generates alpha.",
     "Matches “mirror loses less, not profits.”", TEAL),
    ("Dorfleitner & Scheckenbach (2022)\nJ. of Risk Finance",
     "Higher trading activity is linked to lower returns — overtrading hurts.",
     "Our mirrored wallets are hyperactive whales.", AMBER),
    ("Grundke & Wittke (2024)\nEuro. J. of Finance",
     "Copy-traders herd — but so do mutual funds; social trading is not automatically riskier.",
     "Context for the herding our copy rule induces.", RED),
    ("Kalshi (2018; 2026)\nresearch.kalshi.com",
     "Prediction prices beat Wall St. consensus (−40.1% MAE) and detect voter mobilization "
     "within minutes.",
     "Why prediction-market traders are worth testing at all.", BLUE),
]
y = top + Inches(0.02)
rh = Inches(0.88); gap = Inches(0.11)
for tag, finding, conn, color in lit:
    rrect(s, Inches(0.55), y, Inches(12.2), rh, CARD, radius=0.07)
    rect(s, Inches(0.55), y, Inches(0.10), rh, color)
    tl = tag.split("\n")
    tag_runs = [(tl[0], 12, color, True)] + [(ln, 9.5, SLATE, False) for ln in tl[1:]]
    text(s, Inches(0.78), y + Inches(0.10), Inches(2.95), rh - Inches(0.2),
         tag_runs, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    text(s, Inches(3.95), y + Inches(0.09), Inches(5.55), rh - Inches(0.18),
         [(finding, 11.5, SLATE, False)], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.02)
    text(s, Inches(9.7), y + Inches(0.09), Inches(2.95), rh - Inches(0.18),
         [("→ " + conn, 11, INK, True)], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.02)
    y += rh + gap

# =====================================================================
# 17 — LIMITATIONS
# =====================================================================
s, top = content("Limitations", "Honesty")
lims = [("Wallet ≠ person", "bots, market-makers, copy-traders"),
        ("Exploratory selection", "Study 2 leaderboard rule is survivorship-prone"),
        ("Tiny effective N", "3–4 active wallets, a single session"),
        ("Short window", "mostly same-day live-sports markets"),
        ("Gross of spread", "costs only hurt the long mirror leg"),
        ("SELL ambiguity", "open vs. close conflated (83/5,854)"),
        ("Mark granularity", "at-or-before the horizon timestamp"),
        ("Modeled depth", "no free historical order-book depth (R6)")]
col_w = Inches(6.0); xa, xb = Inches(0.55), Inches(6.75)
y = top + Inches(0.1)
for i, (h, d) in enumerate(lims):
    x = xa if i % 2 == 0 else xb
    if i % 2 == 0 and i > 0:
        y += Inches(1.18)
    rrect(s, x, y, col_w, Inches(1.0), CARD2, radius=0.08, line=RGBColor(0xE0,0xE7,0xEF), line_w=1.0)
    rect(s, x, y, Inches(0.10), Inches(1.0), MUTE)
    text(s, x + Inches(0.28), y + Inches(0.14), col_w - Inches(0.5), Inches(0.75),
         [{"runs": [(h, 14.5, INK, True)], "space_after": 3},
          {"runs": [(d, 12, SLATE, False)], "space_after": 0}], line_spacing=1.0)

# =====================================================================
# 18 — CONCLUSION
# =====================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, INK)
rect(s, 0, 0, Inches(0.16), SH, TEAL)
tri(s, SW - Inches(3.0), SH - Inches(3.0), Inches(3.0), Inches(3.0), INK2, flip=True)
text(s, Inches(0.9), Inches(0.85), Inches(11.5), Inches(0.4),
     [("CONCLUSION", 14, TEAL, True)])
text(s, Inches(0.9), Inches(1.25), Inches(11.5), Inches(1.4),
     [("No robust, exploitable mirror-trading edge.", 36, WHITE, True)], line_spacing=1.0)
rect(s, Inches(0.92), Inches(2.7), Inches(2.2), Pt(3.5), TEAL)
stat_card(s, Inches(0.9), Inches(3.15), Inches(5.8), Inches(1.25), "0 / 151",
          "historical luck-filter survivors", accent=RED, big_size=26, fill=INK2)
stat_card(s, Inches(6.95), Inches(3.15), Inches(5.45), Inches(1.25), "≈ 0",
          "forward edge, clustered by market", accent=TEAL, big_size=26, fill=INK2)
text(s, Inches(0.9), Inches(4.75), Inches(11.5), Inches(1.6),
     [("These markets price observable trading information efficiently enough that neither "
       "past-calibration nor copy-the-winner selection beats the favorite after honest accounting.",
       18, CLOUD, False)], line_spacing=1.18)
text(s, Inches(0.9), Inches(6.45), Inches(11.5), Inches(0.6),
     [("A clean null — and, by the standard we set, a successful result.", 17, TEAL, True)])

# =====================================================================
# 20 — WORKS CITED
# =====================================================================
s, top = content("Works Cited", "References")
refs = [
    "Akey, Pat, et al. “Who Wins and Who Loses in Prediction Markets? Evidence from Polymarket.” "
    "SSRN, 18 Mar. 2026, ssrn.com/abstract=6443103.",
    "Dorfleitner, Gregor, and Isabel Scheckenbach. “Trading Activity on Social Trading Platforms – "
    "a Behavioral Approach.” The Journal of Risk Finance, vol. 23, no. 1, Jan. 2022, pp. 32–54.",
    "Grundke, Peter, and Gerrit Wittke. “Social Trading Platforms vs. Mutual Funds: Herding "
    "Tendencies and Portfolio Risks.” The European Journal of Finance, vol. 31, no. 7, Dec. 2024, "
    "pp. 827–49.",
    "“Kalshi.” Kalshi.com, 2018, research.kalshi.com/articles/crisis-alpha.",
    "“Kalshi Research.” Kalshi.com, 2026, research.kalshi.com/articles/mamdani-primary-victory.",
    "Lu, Juye Shirley. “To Mirror or Not to Mirror: Modeling Relationships in Social Trading.” MIT, "
    "2017, dspace.mit.edu/entities/publication/d04baf10-ccf7-4358-9aaf-6716565782e9.",
]
paras = [{"runs": [(r, 13.5, INK, False)], "space_after": 13} for r in refs]
text(s, Inches(0.6), top + Inches(0.15), Inches(12.1), Inches(4.7), paras, line_spacing=1.08)
text(s, Inches(0.6), SH - Inches(0.95), Inches(12.1), Inches(0.4),
     [("In-text evidence is cited on the motivation, Study 1, and related-work slides; full notes "
       "in the research organizer.", 10.5, MUTE, False)])

saved = None
for cand in ["polymirror_symposium_final.pptx", "polymirror_symposium_final_v2.pptx",
             "polymirror_symposium_final_v3.pptx", "polymirror_symposium_final_v4.pptx"]:
    try:
        prs.save(str(HERE / cand)); saved = cand; break
    except PermissionError:
        continue
if saved is None:
    raise SystemExit("All candidate filenames are locked — close the open PowerPoint files.")
print("wrote", HERE / saved, "with", len(prs.slides._sldIdLst), "slides")
