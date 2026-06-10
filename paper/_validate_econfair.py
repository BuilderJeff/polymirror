"""Validate paper/polymirror_econfair.pptx: slide count, banned jargon,
and that the dynamic numbers match data/forward at this moment."""
from __future__ import annotations

import csv
import json
import re
import sys
from pathlib import Path

from pptx import Presentation

HERE = Path(__file__).resolve().parent
ROOT = HERE.parent
FWD = ROOT / "data" / "forward"

path = HERE / "polymirror_econfair.pptx"
prs = Presentation(str(path))

# ---- extract all text ------------------------------------------------------
texts = []
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    texts.append(r.text)
ALL = " ".join(texts)
errors = []

# ---- slide count ------------------------------------------------------------
n_slides = len(prs.slides._sldIdLst)
if n_slides != 16:
    errors.append(f"slide count {n_slides} != 16")

# ---- banned jargon -----------------------------------------------------------
low = ALL.lower()
for banned in ["bootstrap", "benjamini", "null hypothesis", "p-value", "pvalue",
               "clustered"]:
    if banned in low:
        errors.append(f"banned word present: {banned!r}")
for banned in ["fdr", "clob"]:
    if re.search(rf"\b{banned}\b", low):
        errors.append(f"banned word present: {banned!r}")
n_brier = len(re.findall(r"brier", low))
if n_brier > 1:
    errors.append(f"'Brier' appears {n_brier} times (max 1)")
if n_brier == 1 and \
        "brier score (a score for how close your bets land to what actually happened)" not in low:
    errors.append("'Brier' present but not immediately defined in parentheses")

# ---- dynamic numbers vs data -------------------------------------------------
with open(FWD / "experiment.json", encoding="utf-8") as f:
    exp = json.load(f)
plist = list(exp["positions"].values()) if isinstance(exp["positions"], dict) \
    else list(exp["positions"])
n_pos = len(plist)
n_markets = len({p["condition_id"] for p in plist})
n_wallets = len({p["wallet"] for p in plist})
n_marks = sum(len(p.get("marks", {})) for p in plist)
n_buys = sum(1 for p in plist if str(p.get("side", "")).upper() == "BUY")
n_sells = sum(1 for p in plist if str(p.get("side", "")).upper() == "SELL")

rows = []
with open(FWD / "edge_curve.csv", encoding="utf-8") as f:
    for row in csv.DictReader(f):
        if not row.get("strat") or not row.get("edge"):
            continue
        if int(float(row["n_positions"] or 0)) == 0:
            continue
        rows.append((int(row["horizon_h"]), float(row["edge"])))
rows.sort()
edge_1 = rows[0][1]
later_neg = sum(1 for _, e in rows[1:] if e < 0)
if edge_1 > 0 and later_neg > 0:
    headline = "Copying beat the favorite early — then fell behind."
elif edge_1 > 0:
    headline = "Copying held a small early edge over the favorite."
else:
    headline = "Copying never reliably beat just betting the favorite."

checks = [
    (f"{n_pos:,} copied trades", "slide 4 n_pos"),
    (f"{n_pos:,} mirrored trades", "slide 15 n_pos"),
    (f"{n_buys:,} buys and {n_sells:,} sells", "slide 8 buy/sell split"),
    (f"{n_marks:,} hourly price checks", "slide 9 n_marks"),
    (f"Only {n_wallets} of the 10", "slide 14 active wallets"),
    (headline, "slide 10 computed headline"),
]
for needle, label in checks:
    if needle not in ALL:
        errors.append(f"missing dynamic text ({label}): {needle!r}")

# n_markets / n_wallets stat blocks (standalone numbers)
for needle, label in [(str(n_markets), "n_markets stat"),
                      (str(n_wallets), "n_wallets stat")]:
    if needle not in texts:
        errors.append(f"missing stat block run ({label}): {needle!r}")

print(f"slides={n_slides}  n_pos={n_pos}  n_markets={n_markets}  "
      f"n_wallets_active={n_wallets}  n_marks={n_marks}  buys={n_buys}  "
      f"sells={n_sells}  edge_1h={edge_1:+.4f}  later_neg={later_neg}/{len(rows)-1}")
if errors:
    print("FAIL")
    for e in errors:
        print(" -", e)
    sys.exit(1)
print("PASS: 16 slides, no banned jargon, dynamic numbers match data/forward")
